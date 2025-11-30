"""
PPTX Processing Module - Extraction complète des fichiers PowerPoint

Fonctionnalités:
- Extraction du texte des slides (titres, corps)
- Extraction des images pour OCR LLM
- Extraction des SmartArt (texte des formes)
- Extraction des Charts (titre + légendes + OCR image)
- Extraction des pièces jointes embarquées (OLE objects)
- Gestion des fichiers accompagnants

v1.0 - Support PPTX complet avec OCR
"""

import os
import io
import re
import logging
import tempfile
import zipfile
import base64
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# Tenter d'importer python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("[PPTX] python-pptx not installed. Install with: pip install python-pptx")

# Tenter d'importer Pillow pour l'export d'images
try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False
    logger.warning("[PPTX] Pillow not installed. Image extraction limited.")


# =============================================================================
#  DATA CLASSES
# =============================================================================

@dataclass
class SlideContent:
    """Contenu extrait d'une slide."""
    slide_number: int
    title: str = ""
    body_text: str = ""
    shapes_text: List[str] = field(default_factory=list)
    smartart_text: List[str] = field(default_factory=list)
    table_text: List[str] = field(default_factory=list)
    chart_info: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)  # {data, format, description}


@dataclass
class PPTXExtractionResult:
    """Résultat complet de l'extraction PPTX."""
    filename: str
    total_slides: int
    slides: List[SlideContent]
    embedded_files: List[Dict[str, Any]]  # {filename, type, data}
    full_text: str  # Texte concaténé de toutes les slides
    images_for_ocr: List[Dict[str, Any]]  # Images à traiter par OCR
    metadata: Dict[str, Any]


# =============================================================================
#  PPTX EXTRACTOR
# =============================================================================

class PPTXExtractor:
    """
    Extracteur complet pour fichiers PowerPoint (.pptx).

    Extrait:
    - Texte des slides (titres, corps, formes)
    - SmartArt (texte des formes)
    - Charts (titre + légendes)
    - Images (pour OCR)
    - Tableaux
    - Pièces jointes embarquées (OLE objects)
    """

    def __init__(
        self,
        extract_images: bool = True,
        extract_embedded: bool = True,
        extract_charts: bool = True,
        extract_tables: bool = True,
        log: Optional[logging.Logger] = None
    ):
        """
        Args:
            extract_images: Extraire les images pour OCR
            extract_embedded: Extraire les fichiers embarqués
            extract_charts: Extraire les infos des graphiques
            extract_tables: Extraire le contenu des tableaux
            log: Logger optionnel
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx required for PPTX processing. "
                "Install with: pip install python-pptx"
            )

        self.extract_images = extract_images
        self.extract_embedded = extract_embedded
        self.extract_charts = extract_charts
        self.extract_tables = extract_tables
        self._log = log or logger

    def _extract_text_from_shape(self, shape) -> str:
        """Extrait le texte d'une forme."""
        text_parts = []

        if hasattr(shape, "text") and shape.text:
            text_parts.append(shape.text.strip())

        # Pour les formes avec text_frame
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                for run in paragraph.runs:
                    if run.text:
                        para_text += run.text
                if para_text.strip():
                    text_parts.append(para_text.strip())

        return "\n".join(text_parts)

    def _extract_table_text(self, table) -> List[str]:
        """Extrait le texte d'un tableau."""
        rows_text = []

        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cells_text.append(cell_text)
            if any(cells_text):
                rows_text.append(" | ".join(cells_text))

        return rows_text

    def _extract_chart_info(self, chart) -> Dict[str, Any]:
        """Extrait les informations d'un graphique."""
        chart_info = {
            "type": "chart",
            "title": "",
            "categories": [],
            "series_names": [],
            "description": "",
        }

        try:
            # Titre du chart
            if hasattr(chart, "chart_title") and chart.chart_title:
                if hasattr(chart.chart_title, "text_frame"):
                    chart_info["title"] = chart.chart_title.text_frame.text.strip()

            # Type de chart
            if hasattr(chart, "chart_type"):
                chart_info["chart_type"] = str(chart.chart_type)

            # Catégories (axe X)
            if hasattr(chart, "plots") and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, "categories"):
                        chart_info["categories"] = [str(cat) for cat in plot.categories][:10]

            # Noms des séries
            if hasattr(chart, "series"):
                for series in chart.series:
                    if hasattr(series, "name") and series.name:
                        chart_info["series_names"].append(str(series.name))

            # Construire une description textuelle
            desc_parts = []
            if chart_info["title"]:
                desc_parts.append(f"Graphique: {chart_info['title']}")
            if chart_info["categories"]:
                desc_parts.append(f"Catégories: {', '.join(chart_info['categories'][:5])}")
            if chart_info["series_names"]:
                desc_parts.append(f"Séries: {', '.join(chart_info['series_names'][:5])}")

            chart_info["description"] = " | ".join(desc_parts) if desc_parts else "Graphique sans titre"

        except Exception as e:
            self._log.warning(f"[PPTX] Error extracting chart info: {e}")
            chart_info["description"] = "Graphique (extraction partielle)"

        return chart_info

    def _extract_smartart_text(self, shape) -> List[str]:
        """Extrait le texte d'un SmartArt."""
        texts = []

        try:
            # SmartArt est composé de formes groupées
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    text = self._extract_text_from_shape(sub_shape)
                    if text:
                        texts.append(text)

            # Essayer aussi l'extraction directe
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())

        except Exception as e:
            self._log.debug(f"[PPTX] SmartArt extraction: {e}")

        return texts

    def _extract_image(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """Extrait une image d'une forme."""
        if not self.extract_images:
            return None

        try:
            if hasattr(shape, "image"):
                image = shape.image
                image_data = image.blob
                image_format = image.ext  # png, jpg, etc.

                return {
                    "slide_number": slide_number,
                    "format": image_format,
                    "data": image_data,
                    "size_bytes": len(image_data),
                    "content_type": image.content_type,
                }
        except Exception as e:
            self._log.debug(f"[PPTX] Image extraction failed: {e}")

        return None

    def _process_slide(self, slide, slide_number: int) -> SlideContent:
        """Traite une slide complète."""
        content = SlideContent(slide_number=slide_number)

        # Titre de la slide
        if slide.shapes.title:
            content.title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            try:
                shape_type = shape.shape_type if hasattr(shape, "shape_type") else None

                # Placeholder (titre, sous-titre, corps)
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    text = self._extract_text_from_shape(shape)
                    if text and text != content.title:
                        content.body_text += text + "\n"

                # Tableau
                elif hasattr(shape, "has_table") and shape.has_table:
                    if self.extract_tables:
                        table_text = self._extract_table_text(shape.table)
                        content.table_text.extend(table_text)

                # Graphique (Chart)
                elif hasattr(shape, "has_chart") and shape.has_chart:
                    if self.extract_charts:
                        chart_info = self._extract_chart_info(shape.chart)
                        content.chart_info.append(chart_info)

                # SmartArt (diagramme)
                elif shape_type == MSO_SHAPE_TYPE.GROUP:
                    smartart_texts = self._extract_smartart_text(shape)
                    content.smartart_text.extend(smartart_texts)

                # Image
                elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = self._extract_image(shape, slide_number)
                    if image_data:
                        content.images.append(image_data)

                # Autre forme avec texte
                else:
                    text = self._extract_text_from_shape(shape)
                    if text:
                        content.shapes_text.append(text)

            except Exception as e:
                self._log.debug(f"[PPTX] Shape processing error on slide {slide_number}: {e}")

        return content

    def _extract_embedded_files(self, pptx_path: str) -> List[Dict[str, Any]]:
        """Extrait les fichiers embarqués (OLE objects) du PPTX."""
        if not self.extract_embedded:
            return []

        embedded_files = []

        try:
            # Le PPTX est un fichier ZIP
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                # Chercher dans ppt/embeddings/
                for name in zf.namelist():
                    if name.startswith("ppt/embeddings/"):
                        filename = os.path.basename(name)
                        if filename:
                            data = zf.read(name)
                            ext = os.path.splitext(filename)[1].lower()

                            embedded_files.append({
                                "filename": filename,
                                "original_path": name,
                                "extension": ext,
                                "data": data,
                                "size_bytes": len(data),
                                "type": self._guess_embedded_type(ext),
                            })
                            self._log.info(f"[PPTX] Found embedded file: {filename} ({len(data)} bytes)")

        except Exception as e:
            self._log.warning(f"[PPTX] Error extracting embedded files: {e}")

        return embedded_files

    def _guess_embedded_type(self, ext: str) -> str:
        """Devine le type d'un fichier embarqué."""
        ext_types = {
            ".pdf": "pdf",
            ".doc": "word",
            ".docx": "word",
            ".xls": "excel",
            ".xlsx": "excel",
            ".png": "image",
            ".jpg": "image",
            ".jpeg": "image",
            ".gif": "image",
            ".bmp": "image",
            ".ole": "ole_object",
            ".bin": "binary",
        }
        return ext_types.get(ext, "unknown")

    def _extract_metadata(self, prs: "Presentation") -> Dict[str, Any]:
        """Extrait les métadonnées du document."""
        metadata = {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width.inches if hasattr(prs.slide_width, "inches") else None,
            "slide_height": prs.slide_height.inches if hasattr(prs.slide_height, "inches") else None,
        }

        # Core properties
        try:
            core_props = prs.core_properties
            if core_props:
                metadata["title"] = core_props.title or ""
                metadata["author"] = core_props.author or ""
                metadata["subject"] = core_props.subject or ""
                metadata["created"] = str(core_props.created) if core_props.created else ""
                metadata["modified"] = str(core_props.modified) if core_props.modified else ""
        except Exception as e:
            self._log.debug(f"[PPTX] Metadata extraction: {e}")

        return metadata

    def extract(self, pptx_path: str) -> PPTXExtractionResult:
        """
        Extrait le contenu complet d'un fichier PPTX.

        Args:
            pptx_path: Chemin vers le fichier .pptx

        Returns:
            PPTXExtractionResult avec tout le contenu extrait
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

        self._log.info(f"[PPTX] Extracting: {pptx_path}")

        # Ouvrir le fichier PPTX
        prs = Presentation(pptx_path)

        # Extraire les métadonnées
        metadata = self._extract_metadata(prs)

        # Traiter chaque slide
        slides_content = []
        all_images = []

        for i, slide in enumerate(prs.slides, start=1):
            self._log.debug(f"[PPTX] Processing slide {i}/{len(prs.slides)}")
            slide_content = self._process_slide(slide, i)
            slides_content.append(slide_content)

            # Collecter les images pour OCR
            all_images.extend(slide_content.images)

        # Extraire les fichiers embarqués
        embedded_files = self._extract_embedded_files(pptx_path)

        # Construire le texte complet
        full_text = self._build_full_text(slides_content)

        result = PPTXExtractionResult(
            filename=os.path.basename(pptx_path),
            total_slides=len(prs.slides),
            slides=slides_content,
            embedded_files=embedded_files,
            full_text=full_text,
            images_for_ocr=all_images,
            metadata=metadata,
        )

        self._log.info(
            f"[PPTX] Extraction complete: {result.total_slides} slides, "
            f"{len(all_images)} images, {len(embedded_files)} embedded files"
        )

        return result

    def _build_full_text(self, slides: List[SlideContent]) -> str:
        """Construit le texte complet à partir des slides."""
        text_parts = []

        for slide in slides:
            slide_text = []
            slide_text.append(f"=== SLIDE {slide.slide_number} ===")

            if slide.title:
                slide_text.append(f"Titre: {slide.title}")

            if slide.body_text.strip():
                slide_text.append(slide.body_text.strip())

            if slide.shapes_text:
                slide_text.append("\n".join(slide.shapes_text))

            if slide.smartart_text:
                slide_text.append("SmartArt: " + " | ".join(slide.smartart_text))

            if slide.table_text:
                slide_text.append("Tableau:\n" + "\n".join(slide.table_text))

            if slide.chart_info:
                for chart in slide.chart_info:
                    slide_text.append(chart.get("description", "Graphique"))

            text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)


# =============================================================================
#  OCR INTEGRATION
# =============================================================================

def process_pptx_images_with_ocr(
    images: List[Dict[str, Any]],
    ocr_func,
    log: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """
    Traite les images extraites du PPTX avec OCR.

    Args:
        images: Liste d'images extraites {data, format, slide_number}
        ocr_func: Fonction OCR à appeler (image_data -> text)
        log: Logger optionnel

    Returns:
        Liste d'images avec texte OCR ajouté
    """
    _log = log or logger

    for i, img in enumerate(images):
        try:
            _log.info(f"[PPTX-OCR] Processing image {i+1}/{len(images)} from slide {img.get('slide_number')}")

            # Appeler la fonction OCR
            ocr_text = ocr_func(img["data"])

            if ocr_text:
                img["ocr_text"] = ocr_text
                _log.debug(f"[PPTX-OCR] Extracted {len(ocr_text)} chars from image")
            else:
                img["ocr_text"] = ""

        except Exception as e:
            _log.warning(f"[PPTX-OCR] OCR failed for image {i+1}: {e}")
            img["ocr_text"] = ""
            img["ocr_error"] = str(e)

    return images


def process_pptx_charts_with_ocr(
    pptx_path: str,
    slides: List[SlideContent],
    ocr_func,
    log: Optional[logging.Logger] = None
) -> List[SlideContent]:
    """
    Capture les graphiques en images et applique l'OCR.

    Note: Nécessite une conversion externe (ex: LibreOffice) pour
    exporter les slides en images. Cette fonction est un placeholder
    pour l'intégration future.

    Args:
        pptx_path: Chemin du fichier PPTX
        slides: Contenu des slides
        ocr_func: Fonction OCR
        log: Logger

    Returns:
        Slides avec OCR des charts ajouté
    """
    _log = log or logger

    # Pour l'instant, on utilise la description textuelle des charts
    # L'OCR sur les images de charts nécessiterait une conversion externe
    _log.info("[PPTX-OCR] Chart OCR: using text descriptions (image export not implemented)")

    return slides


# =============================================================================
#  EMBEDDED FILES PROCESSING
# =============================================================================

def process_embedded_files(
    embedded_files: List[Dict[str, Any]],
    pdf_processor=None,
    image_ocr_func=None,
    log: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """
    Traite les fichiers embarqués extraits du PPTX.

    Args:
        embedded_files: Liste des fichiers embarqués
        pdf_processor: Fonction pour traiter les PDFs (pdf_data -> text)
        image_ocr_func: Fonction OCR pour les images
        log: Logger

    Returns:
        Fichiers embarqués avec texte extrait
    """
    _log = log or logger

    for f in embedded_files:
        try:
            file_type = f.get("type", "unknown")
            _log.info(f"[PPTX-EMBED] Processing {f['filename']} (type={file_type})")

            if file_type == "pdf" and pdf_processor:
                # Sauvegarder temporairement le PDF
                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                    tmp.write(f["data"])
                    tmp_path = tmp.name

                try:
                    f["extracted_text"] = pdf_processor(tmp_path)
                finally:
                    os.unlink(tmp_path)

            elif file_type == "image" and image_ocr_func:
                f["extracted_text"] = image_ocr_func(f["data"])

            else:
                f["extracted_text"] = ""
                _log.debug(f"[PPTX-EMBED] No processor for type: {file_type}")

        except Exception as e:
            _log.warning(f"[PPTX-EMBED] Processing failed for {f['filename']}: {e}")
            f["extracted_text"] = ""
            f["processing_error"] = str(e)

    return embedded_files


# =============================================================================
#  UTILITY FUNCTIONS
# =============================================================================

def extract_pptx(
    pptx_path: str,
    extract_images: bool = True,
    extract_embedded: bool = True,
    log: Optional[logging.Logger] = None
) -> PPTXExtractionResult:
    """
    Fonction utilitaire pour extraire un PPTX.

    Args:
        pptx_path: Chemin du fichier PPTX
        extract_images: Extraire les images
        extract_embedded: Extraire les fichiers embarqués
        log: Logger

    Returns:
        Résultat de l'extraction
    """
    extractor = PPTXExtractor(
        extract_images=extract_images,
        extract_embedded=extract_embedded,
        log=log
    )
    return extractor.extract(pptx_path)


def pptx_to_text(pptx_path: str, include_metadata: bool = True) -> str:
    """
    Convertit un PPTX en texte simple.

    Args:
        pptx_path: Chemin du fichier PPTX
        include_metadata: Inclure les métadonnées

    Returns:
        Texte extrait du PPTX
    """
    result = extract_pptx(pptx_path, extract_images=False, extract_embedded=False)

    if include_metadata:
        meta_str = f"Document: {result.filename}\n"
        meta_str += f"Slides: {result.total_slides}\n"
        if result.metadata.get("title"):
            meta_str += f"Title: {result.metadata['title']}\n"
        if result.metadata.get("author"):
            meta_str += f"Author: {result.metadata['author']}\n"
        meta_str += "\n"
        return meta_str + result.full_text

    return result.full_text


def find_accompanying_files(
    pptx_path: str,
    extensions: Optional[List[str]] = None
) -> List[str]:
    """
    Trouve les fichiers accompagnant un PPTX (dans le même dossier).

    Args:
        pptx_path: Chemin du fichier PPTX
        extensions: Extensions à rechercher (défaut: images et PDFs)

    Returns:
        Liste des chemins des fichiers accompagnants
    """
    if extensions is None:
        extensions = [".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]

    pptx_dir = os.path.dirname(pptx_path) or "."
    pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]

    accompanying = []

    for filename in os.listdir(pptx_dir):
        name, ext = os.path.splitext(filename)
        ext_lower = ext.lower()

        # Chercher les fichiers avec le même préfixe ou dans un sous-dossier
        if ext_lower in extensions:
            # Même préfixe (ex: presentation_image1.png)
            if name.startswith(pptx_name):
                accompanying.append(os.path.join(pptx_dir, filename))

    # Chercher aussi dans un sous-dossier du même nom
    sub_folder = os.path.join(pptx_dir, pptx_name)
    if os.path.isdir(sub_folder):
        for filename in os.listdir(sub_folder):
            _, ext = os.path.splitext(filename)
            if ext.lower() in extensions:
                accompanying.append(os.path.join(sub_folder, filename))

    return accompanying


def process_pptx_with_attachments(
    pptx_path: str,
    ocr_func=None,
    pdf_processor=None,
    log: Optional[logging.Logger] = None
) -> Dict[str, Any]:
    """
    Traitement complet d'un PPTX avec toutes les pièces jointes.

    Combine:
    - Extraction du texte des slides
    - OCR des images dans les slides
    - Extraction et traitement des fichiers embarqués
    - Traitement des fichiers accompagnants

    Args:
        pptx_path: Chemin du fichier PPTX
        ocr_func: Fonction OCR pour les images (image_data -> text)
        pdf_processor: Fonction pour traiter les PDFs (pdf_path -> text)
        log: Logger

    Returns:
        Dict avec:
        - full_text: Texte complet (slides + OCR + embedded)
        - slides: Contenu des slides
        - embedded_files: Fichiers embarqués traités
        - accompanying_files: Fichiers accompagnants traités
        - metadata: Métadonnées du document
    """
    _log = log or logger
    _log.info(f"[PPTX] Full processing: {pptx_path}")

    # 1. Extraction de base
    result = extract_pptx(
        pptx_path,
        extract_images=True,
        extract_embedded=True,
        log=_log
    )

    # 2. OCR des images si fonction fournie
    if ocr_func and result.images_for_ocr:
        _log.info(f"[PPTX] Processing {len(result.images_for_ocr)} images with OCR")
        result.images_for_ocr = process_pptx_images_with_ocr(
            result.images_for_ocr,
            ocr_func,
            log=_log
        )

    # 3. Traitement des fichiers embarqués
    if result.embedded_files:
        _log.info(f"[PPTX] Processing {len(result.embedded_files)} embedded files")
        result.embedded_files = process_embedded_files(
            result.embedded_files,
            pdf_processor=pdf_processor,
            image_ocr_func=ocr_func,
            log=_log
        )

    # 4. Fichiers accompagnants
    accompanying = find_accompanying_files(pptx_path)
    accompanying_content = []

    for acc_path in accompanying:
        _log.info(f"[PPTX] Processing accompanying file: {acc_path}")
        try:
            ext = os.path.splitext(acc_path)[1].lower()

            if ext == ".pdf" and pdf_processor:
                text = pdf_processor(acc_path)
            elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp") and ocr_func:
                with open(acc_path, "rb") as f:
                    img_data = f.read()
                text = ocr_func(img_data)
            else:
                text = ""

            accompanying_content.append({
                "path": acc_path,
                "filename": os.path.basename(acc_path),
                "type": ext[1:],
                "extracted_text": text,
            })

        except Exception as e:
            _log.warning(f"[PPTX] Failed to process accompanying file {acc_path}: {e}")

    # 5. Construire le texte complet enrichi
    full_text_parts = [result.full_text]

    # Ajouter le texte OCR des images
    for img in result.images_for_ocr:
        if img.get("ocr_text"):
            full_text_parts.append(
                f"\n[Image Slide {img['slide_number']} - OCR]\n{img['ocr_text']}"
            )

    # Ajouter le texte des fichiers embarqués
    for emb in result.embedded_files:
        if emb.get("extracted_text"):
            full_text_parts.append(
                f"\n[Fichier embarqué: {emb['filename']}]\n{emb['extracted_text']}"
            )

    # Ajouter le texte des fichiers accompagnants
    for acc in accompanying_content:
        if acc.get("extracted_text"):
            full_text_parts.append(
                f"\n[Fichier accompagnant: {acc['filename']}]\n{acc['extracted_text']}"
            )

    return {
        "full_text": "\n\n".join(full_text_parts),
        "slides": result.slides,
        "embedded_files": result.embedded_files,
        "accompanying_files": accompanying_content,
        "images_for_ocr": result.images_for_ocr,
        "metadata": result.metadata,
        "total_slides": result.total_slides,
        "filename": result.filename,
    }


# =============================================================================
#  CHECK AVAILABILITY
# =============================================================================

def is_pptx_processing_available() -> bool:
    """Vérifie si le traitement PPTX est disponible."""
    return PPTX_AVAILABLE
